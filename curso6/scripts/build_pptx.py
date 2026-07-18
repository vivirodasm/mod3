"""Generate dia1.pptx / dia2.pptx from each day's ppt_contenido.md."""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

CURSO6 = Path(__file__).resolve().parents[1]

TITLE_COLOR = RGBColor(0x1B, 0x2A, 0x4A)
BODY_COLOR = RGBColor(0x2C, 0x3E, 0x50)
ACCENT_COLOR = RGBColor(0x0E, 0x7C, 0x7B)
MUTED_COLOR = RGBColor(0x7F, 0x8C, 0x8D)


def parse_slides(markdown: str) -> list[dict[str, object]]:
    slides: list[dict[str, object]] = []
    chunks = re.split(r"\n### (Slide \d+ —[^\n]+)\n", markdown)
    for i in range(1, len(chunks), 2):
        heading = chunks[i].strip()
        body = chunks[i + 1] if i + 1 < len(chunks) else ""
        title_match = re.match(r"Slide (\d+) — (.+?)(?:\s*\(\d+ min\))?$", heading)
        if not title_match:
            continue
        number = int(title_match.group(1))
        title = title_match.group(2).strip()

        bullets: list[str] = []
        code_lines: list[str] = []
        in_code = False
        for line in body.splitlines():
            if line.startswith("**Título:**"):
                title = line.replace("**Título:**", "").strip()
                continue
            if line.strip().startswith("```"):
                in_code = not in_code
                continue
            if in_code:
                if line.strip():
                    code_lines.append(line.rstrip())
                continue
            if line.startswith("---") or line.startswith("## "):
                break
            if line.startswith("- "):
                bullets.append(line[2:].strip())
            elif line.startswith("| ") and not re.match(r"^\|\s*-+", line):
                cells = [c.strip() for c in line.strip("|").split("|")]
                if cells and cells[0].lower() not in {
                    "herramienta",
                    "métrica",
                    "plataforma",
                    "si duele…",
                }:
                    bullets.append(" — ".join(cells))
        if code_lines:
            bullets.append("Código: " + " | ".join(code_lines[:5]))
        slides.append({"number": number, "title": title, "bullets": bullets[:8]})
    return slides


def style_run(run, text: str, *, size: int, bold: bool = False, color=BODY_COLOR) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def clean_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    return text.replace("`", "")


def add_slide(
    prs: Presentation, number: int, title: str, bullets: list[str], day_label: str
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    title_box.text_frame.word_wrap = True
    style_run(p.add_run(), title, size=28, bold=True, color=TITLE_COLOR)

    meta = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.3))
    style_run(
        meta.text_frame.paragraphs[0].add_run(),
        f"Certificación 6 · {day_label} · Slide {number}",
        size=12,
        color=ACCENT_COLOR,
    )

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    items = bullets or ["Ver guía del instructor para el relato completo."]
    for idx, bullet in enumerate(items):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.level = 0
        para.space_after = Pt(10)
        style_run(para.add_run(), f"• {clean_md(bullet)}", size=18, color=BODY_COLOR)

    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3))
    fp = footer.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    style_run(
        fp.add_run(),
        "Madurez, Gobernanza y Estrategia de Calidad",
        size=10,
        color=MUTED_COLOR,
    )


def build_day(day_dir: Path) -> Path:
    md_path = day_dir / "ppt_contenido.md"
    out_path = day_dir / f"{day_dir.name}.pptx"
    if not md_path.exists():
        raise SystemExit(f"No existe {md_path}")

    slides = parse_slides(md_path.read_text(encoding="utf-8"))
    if not slides:
        raise SystemExit(f"No se encontraron slides en {md_path}")

    day_label = "Día 1" if day_dir.name == "dia1" else "Día 2"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide in slides:
        add_slide(
            prs,
            int(slide["number"]),
            str(slide["title"]),
            list(slide["bullets"]),
            day_label,
        )
    prs.save(out_path)
    print(f"Generado {out_path} con {len(slides)} diapositivas")
    return out_path


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "day",
        nargs="?",
        choices=["dia1", "dia2", "all"],
        default="all",
        help="Qué día generar (default: all)",
    )
    args = parser.parse_args()
    days = ["dia1", "dia2"] if args.day == "all" else [args.day]
    for name in days:
        build_day(CURSO6 / name)


if __name__ == "__main__":
    main()
